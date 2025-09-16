"""
Alternative Document Processor - 로컬 Python 라이브러리를 사용한 문서 처리

Docling 서비스가 사용 불가능할 때의 대안 처리기입니다.
python-docx, python-pptx, openpyxl, PyPDF2 등을 사용합니다.
"""

import os
import io
from typing import Dict, Any, Optional, Tuple
import logging

logger = logging.getLogger(__name__)

class AlternativeProcessor:
    """로컬 Python 라이브러리를 사용한 문서 처리기"""
    
    def __init__(self):
        self.available_libraries = self._check_available_libraries()
        
    def _check_available_libraries(self) -> Dict[str, bool]:
        """사용 가능한 라이브러리 확인"""
        libraries = {}
        
        try:
            import docx
            libraries['docx'] = True
        except ImportError:
            libraries['docx'] = False
            
        try:
            from pptx import Presentation
            libraries['pptx'] = True
        except ImportError:
            libraries['pptx'] = False
            
        try:
            import openpyxl
            libraries['openpyxl'] = True
        except ImportError:
            libraries['openpyxl'] = False
            
        try:
            import PyPDF2
            libraries['pypdf2'] = True
        except ImportError:
            libraries['pypdf2'] = False
            
        return libraries
    
    def is_available(self) -> bool:
        """최소 하나의 라이브러리라도 사용 가능한지 확인"""
        return any(self.available_libraries.values())
    
    async def process_document(self, file_content: bytes, filename: str) -> Tuple[bool, Dict[str, Any]]:
        """
        문서를 로컬 라이브러리로 처리
        
        Args:
            file_content: 파일의 바이트 내용
            filename: 파일명 (확장자 포함)
            
        Returns:
            (성공여부, 결과데이터)
        """
        ext = os.path.splitext(filename.lower())[1]
        
        try:
            if ext in ['.docx', '.doc'] and self.available_libraries.get('docx'):
                return await self._process_docx(file_content, filename)
            elif ext in ['.pptx', '.ppt'] and self.available_libraries.get('pptx'):
                return await self._process_pptx(file_content, filename)
            elif ext in ['.xlsx', '.xls'] and self.available_libraries.get('openpyxl'):
                return await self._process_excel(file_content, filename)
            elif ext == '.pdf' and self.available_libraries.get('pypdf2'):
                return await self._process_pdf(file_content, filename)
            else:
                return False, {
                    'method': 'alternative_processor',
                    'error': f'Unsupported file type {ext} or missing library',
                    'status': 'failed'
                }
                
        except Exception as e:
            logger.error(f"Alternative processor error: {e}")
            return False, {
                'method': 'alternative_processor',
                'error': str(e),
                'status': 'failed'
            }
    
    async def _process_docx(self, file_content: bytes, filename: str) -> Tuple[bool, Dict[str, Any]]:
        """DOCX 파일 처리"""
        try:
            import docx
            
            doc_stream = io.BytesIO(file_content)
            doc = docx.Document(doc_stream)
            
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
                    
            content = '\n\n'.join(text_content)
            
            return True, {
                'method': 'alternative_processor',
                'content': content,
                'metadata': {
                    'paragraphs': len(doc.paragraphs),
                    'characters': len(content)
                },
                'status': 'success'
            }
            
        except Exception as e:
            return False, {
                'method': 'alternative_processor',
                'error': f'DOCX processing failed: {str(e)}',
                'status': 'failed'
            }
    
    async def _process_pptx(self, file_content: bytes, filename: str) -> Tuple[bool, Dict[str, Any]]:
        """PPTX 파일 처리"""
        try:
            from pptx import Presentation
            
            ppt_stream = io.BytesIO(file_content)
            prs = Presentation(ppt_stream)
            
            text_content = []
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        
                if slide_text:
                    text_content.append(f"슬라이드 {slide_count}:\n" + '\n'.join(slide_text))
                    
            content = '\n\n'.join(text_content)
            
            return True, {
                'method': 'alternative_processor',
                'content': content,
                'metadata': {
                    'slides': slide_count,
                    'characters': len(content)
                },
                'status': 'success'
            }
            
        except Exception as e:
            return False, {
                'method': 'alternative_processor',
                'error': f'PPTX processing failed: {str(e)}',
                'status': 'failed'
            }
    
    async def _process_excel(self, file_content: bytes, filename: str) -> Tuple[bool, Dict[str, Any]]:
        """Excel 파일 처리"""
        try:
            import openpyxl
            
            excel_stream = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(excel_stream)
            
            text_content = []
            total_rows = 0
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"시트: {sheet_name}"]
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None and str(cell).strip() for cell in row):
                        row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                        sheet_text.append(row_text)
                        total_rows += 1
                        
                if len(sheet_text) > 1:  # 시트 제목 외에 데이터가 있는 경우
                    text_content.append('\n'.join(sheet_text))
                    
            content = '\n\n'.join(text_content)
            
            return True, {
                'method': 'alternative_processor',
                'content': content,
                'metadata': {
                    'sheets': len(workbook.sheetnames),
                    'rows': total_rows,
                    'characters': len(content)
                },
                'status': 'success'
            }
            
        except Exception as e:
            return False, {
                'method': 'alternative_processor',
                'error': f'Excel processing failed: {str(e)}',
                'status': 'failed'
            }
    
    async def _process_pdf(self, file_content: bytes, filename: str) -> Tuple[bool, Dict[str, Any]]:
        """PDF 파일 처리"""
        try:
            import PyPDF2
            
            pdf_stream = io.BytesIO(file_content)
            reader = PyPDF2.PdfReader(pdf_stream)
            
            text_content = []
            page_count = len(reader.pages)
            
            for i, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"페이지 {i}:\n{page_text.strip()}")
                except Exception as e:
                    logger.warning(f"페이지 {i} 추출 실패: {e}")
                    continue
                    
            content = '\n\n'.join(text_content)
            
            return True, {
                'method': 'alternative_processor',
                'content': content,
                'metadata': {
                    'pages': page_count,
                    'characters': len(content)
                },
                'status': 'success'
            }
            
        except Exception as e:
            return False, {
                'method': 'alternative_processor',
                'error': f'PDF processing failed: {str(e)}',
                'status': 'failed'
            }
    
    def is_supported_format(self, filename: str) -> bool:
        """지원되는 파일 형식인지 확인"""
        ext = os.path.splitext(filename.lower())[1]
        
        if ext in ['.docx', '.doc']:
            return self.available_libraries.get('docx', False)
        elif ext in ['.pptx', '.ppt']:
            return self.available_libraries.get('pptx', False)
        elif ext in ['.xlsx', '.xls']:
            return self.available_libraries.get('openpyxl', False)
        elif ext == '.pdf':
            return self.available_libraries.get('pypdf2', False)
        
        return False

# 전역 프로세서 인스턴스
alternative_processor = AlternativeProcessor()