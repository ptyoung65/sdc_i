"""
Multi-format Document Processor for Korean RAG
한글 문서의 다양한 형식을 처리하는 프로세서
"""

import os
import tempfile
import logging
from pathlib import Path
from typing import Dict, Any, Optional, Tuple
import aiohttp
import aiofiles
import asyncio

# Document processing libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """다양한 형식의 문서를 처리하는 클래스"""
    
    def __init__(self, docling_url: str = "http://localhost:5001"):
        self.docling_url = docling_url
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.doc': self._process_doc_via_docling,
            '.docx': self._process_docx,
            '.ppt': self._process_ppt,
            '.pptx': self._process_pptx,
            '.xls': self._process_xls_via_docling,
            '.xlsx': self._process_xlsx,
            '.txt': self._process_txt,
            '.md': self._process_txt,  # Markdown도 텍스트로 처리
        }
    
    async def process_file(self, file_content: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """
        파일을 처리하여 텍스트 추출
        
        Args:
            file_content: 파일 내용
            filename: 파일명
            
        Returns:
            (extracted_text, metadata)
        """
        file_ext = Path(filename).suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {file_ext}")
        
        logger.info(f"Processing file: {filename} ({file_ext})")
        
        # 임시 파일로 저장
        with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        try:
            processor = self.supported_formats[file_ext]
            text, metadata = await processor(temp_path, filename)
            
            # 기본 메타데이터 추가
            metadata.update({
                'filename': filename,
                'file_type': file_ext,
                'file_size': len(file_content)
            })
            
            logger.info(f"Extracted {len(text)} characters from {filename}")
            return text, metadata
            
        finally:
            # 임시 파일 정리
            try:
                os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Failed to cleanup temp file {temp_path}: {e}")
    
    async def _process_txt(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """텍스트/마크다운 파일 처리"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            return content, {'processing_method': 'native_text'}
        except UnicodeDecodeError:
            # UTF-8로 읽기 실패시 다른 인코딩 시도
            encodings = ['cp949', 'euc-kr', 'latin-1']
            for encoding in encodings:
                try:
                    async with aiofiles.open(file_path, 'r', encoding=encoding) as f:
                        content = await f.read()
                    logger.info(f"Successfully decoded {filename} with {encoding}")
                    return content, {'processing_method': 'native_text', 'encoding': encoding}
                except UnicodeDecodeError:
                    continue
            raise ValueError(f"Could not decode text file: {filename}")
    
    async def _process_docx(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """DOCX 파일 처리"""
        if not DOCX_AVAILABLE:
            return await self._process_doc_via_docling(file_path, filename)
        
        try:
            doc = Document(file_path)
            text_parts = []
            
            # 단락별로 텍스트 추출
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            
            # 표 내용도 추출
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            content = '\n\n'.join(text_parts)
            metadata = {
                'processing_method': 'python_docx',
                'paragraphs_count': len(doc.paragraphs),
                'tables_count': len(doc.tables)
            }
            
            return content, metadata
            
        except Exception as e:
            logger.warning(f"Failed to process DOCX with python-docx: {e}")
            return await self._process_doc_via_docling(file_path, filename)
    
    async def _process_ppt(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """PPT 파일 처리 (Docling 우선, 실패시 기본 처리)"""
        # 먼저 Docling 시도
        try:
            return await self._process_ppt_via_docling(file_path, filename)
        except Exception as e:
            logger.warning(f"Docling failed for PPT: {e}")
        
        # python-pptx로 시도 (PPTX와 동일한 라이브러리)
        if PPTX_AVAILABLE:
            try:
                prs = Presentation(file_path)
                text_parts = []
                slide_count = 0
                
                for slide in prs.slides:
                    slide_count += 1
                    slide_texts = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    
                    if slide_texts:
                        text_parts.append(f"슬라이드 {slide_count}:\n" + '\n'.join(slide_texts))
                
                content = '\n\n'.join(text_parts)
                metadata = {
                    'processing_method': 'python_pptx_fallback',
                    'slides_count': slide_count
                }
                
                return content, metadata
                
            except Exception as e:
                logger.warning(f"Failed to process PPT with python-pptx: {e}")
        
        # 최후 수단: 텍스트 파일로 시도 (가장 기본적인 PPT 파일의 경우)
        logger.warning(f"All PPT processing methods failed for {filename}, trying basic text fallback")
        try:
            return await self._process_txt(file_path, filename)
        except Exception as text_error:
            logger.warning(f"Even text fallback failed for PPT {filename}: {text_error}")
        
        # 모든 처리 실패시
        raise ValueError(f"PPT 파일 처리에 실패했습니다: {filename}")

    async def _process_pptx(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """PPTX 파일 처리"""
        if not PPTX_AVAILABLE:
            return await self._process_ppt_via_docling(file_path, filename)
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                slide_texts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    text_parts.append(f"슬라이드 {slide_count}:\n" + '\n'.join(slide_texts))
            
            content = '\n\n'.join(text_parts)
            metadata = {
                'processing_method': 'python_pptx',
                'slides_count': slide_count
            }
            
            return content, metadata
            
        except Exception as e:
            logger.warning(f"Failed to process PPTX with python-pptx: {e}")
            return await self._process_ppt_via_docling(file_path, filename)
    
    async def _process_xlsx(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """XLSX 파일 처리"""
        if not OPENPYXL_AVAILABLE and not PANDAS_AVAILABLE:
            return await self._process_xls_via_docling(file_path, filename)
        
        try:
            if PANDAS_AVAILABLE:
                # pandas 사용
                df = pd.read_excel(file_path, sheet_name=None)  # 모든 시트 읽기
                text_parts = []
                
                for sheet_name, sheet_df in df.items():
                    if not sheet_df.empty:
                        text_parts.append(f"시트: {sheet_name}")
                        # 헤더가 있다면 포함
                        text_parts.append(sheet_df.to_string(index=False))
                
                content = '\n\n'.join(text_parts)
                metadata = {
                    'processing_method': 'pandas',
                    'sheets_count': len(df)
                }
                
            else:
                # openpyxl 사용
                from openpyxl import load_workbook
                wb = load_workbook(file_path)
                text_parts = []
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_texts = [f"시트: {sheet_name}"]
                    
                    for row in sheet.iter_rows(values_only=True):
                        if any(cell for cell in row if cell is not None):
                            row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                            sheet_texts.append(row_text.strip())
                    
                    if len(sheet_texts) > 1:  # 헤더 외에 내용이 있으면
                        text_parts.append('\n'.join(sheet_texts))
                
                content = '\n\n'.join(text_parts)
                metadata = {
                    'processing_method': 'openpyxl',
                    'sheets_count': len(wb.sheetnames)
                }
            
            return content, metadata
            
        except Exception as e:
            logger.warning(f"Failed to process XLSX with local libraries: {e}")
            return await self._process_xls_via_docling(file_path, filename)
    
    async def _process_pdf(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """PDF 파일 처리"""
        # 먼저 Docling 시도
        try:
            return await self._process_via_docling(file_path, filename)
        except Exception as e:
            logger.warning(f"Docling failed for PDF: {e}")
        
        # PyPDF2로 fallback
        if not PYPDF2_AVAILABLE:
            raise ValueError(f"PDF processing not available for {filename}")
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_parts = []
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"페이지 {page_num}:\n{page_text.strip()}")
                
                content = '\n\n'.join(text_parts)
                metadata = {
                    'processing_method': 'pypdf2',
                    'pages_count': len(pdf_reader.pages)
                }
                
                return content, metadata
                
        except Exception as e:
            logger.error(f"Failed to process PDF {filename}: {e}")
            raise ValueError(f"PDF processing failed: {e}")
    
    async def _process_doc_via_docling(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """DOC 파일을 Docling으로 처리"""
        return await self._process_via_docling(file_path, filename)
    
    async def _process_ppt_via_docling(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """PPT 파일을 Docling으로 처리"""
        return await self._process_via_docling(file_path, filename)
    
    async def _process_xls_via_docling(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """XLS 파일을 Docling으로 처리"""
        return await self._process_via_docling(file_path, filename)
    
    async def _process_via_docling(self, file_path: str, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Docling 서비스를 통한 문서 처리"""
        try:
            async with aiohttp.ClientSession() as session:
                # 파일을 multipart로 전송
                with open(file_path, 'rb') as f:
                    data = aiohttp.FormData()
                    data.add_field('file', f, filename=filename)
                    
                    async with session.post(
                        f"{self.docling_url}/convert",
                        data=data,
                        timeout=aiohttp.ClientTimeout(total=60)
                    ) as response:
                        
                        if response.status == 200:
                            result = await response.json()
                            content = result.get('content', '')
                            metadata = {
                                'processing_method': 'docling',
                                'docling_metadata': result.get('metadata', {})
                            }
                            return content, metadata
                        else:
                            error_text = await response.text()
                            raise Exception(f"Docling error {response.status}: {error_text}")
                            
        except Exception as e:
            logger.error(f"Docling processing failed for {filename}: {e}")
            raise ValueError(f"Docling processing failed: {e}")